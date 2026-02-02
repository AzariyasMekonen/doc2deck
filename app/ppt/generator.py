from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
import re
import os

class PPTGenerator:
    def create_presentation(self, notes: str, user_id: int) -> str:
        """Convert notes into a styled PowerPoint presentation"""
        prs = Presentation()
        
        # Parse notes into slides
        slides_content = self._parse_notes_to_slides(notes)
        
        # Create styled title slide
        self._create_title_slide(prs)
        
        # Create content slides with styling
        for slide_content in slides_content:
            self._create_content_slide(prs, slide_content)
        
        # Save presentation
        output_path = f"uploads/presentation_{user_id}.pptx"
        prs.save(output_path)
        
        return output_path
    
    def _create_title_slide(self, prs):
        """Create a styled title slide"""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Style title
        title = slide.shapes.title
        title.text = "Generated Presentation"
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(31, 73, 125)  # Dark blue
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Style subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = "Created with Doc2Deck"
        subtitle_paragraph = subtitle.text_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(24)
        subtitle_paragraph.font.color.rgb = RGBColor(68, 114, 196)  # Light blue
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
    
    def _parse_notes_to_slides(self, notes: str):
        """Parse structured notes into slide content with better organization"""
        slides = []
        current_slide = {"title": "", "content": []}
        
        lines = notes.split('\n')
        current_paragraph = ""
        
        for line in lines:
            line = line.strip()
            
            if not line:
                # End of paragraph - add to current slide if not empty
                if current_paragraph:
                    current_slide["content"].append(current_paragraph)
                    current_paragraph = ""
                continue
            
            # Check for headings
            if line.startswith('# '):
                # Main title - start new slide
                if current_slide["title"] or current_slide["content"]:
                    slides.append(current_slide)
                current_slide = {"title": line[2:], "content": []}
                
            elif line.startswith('## '):
                # Subtitle - start new slide
                if current_slide["title"] or current_slide["content"]:
                    slides.append(current_slide)
                current_slide = {"title": line[3:], "content": []}
                
            elif line.startswith('• ') or line.startswith('- '):
                # Bullet point - add to current slide
                if current_paragraph:
                    current_slide["content"].append(current_paragraph)
                    current_paragraph = ""
                bullet_text = line[2:] if line.startswith('• ') else line[2:]
                current_slide["content"].append(bullet_text)
                
            else:
                # Regular text - add to current paragraph
                if current_paragraph:
                    current_paragraph += " " + line
                else:
                    current_paragraph = line
        
        # Add final paragraph and slide
        if current_paragraph:
            current_slide["content"].append(current_paragraph)
        if current_slide["title"] or current_slide["content"]:
            slides.append(current_slide)
        
        # Balance slide content - split overly long slides
        balanced_slides = []
        for slide in slides:
            if len(slide["content"]) > 6:  # Too much content
                # Split into multiple slides
                content_chunks = [slide["content"][i:i+4] for i in range(0, len(slide["content"]), 4)]
                for i, chunk in enumerate(content_chunks):
                    title_suffix = f" (Part {i+1})" if len(content_chunks) > 1 else ""
                    balanced_slides.append({
                        "title": slide["title"] + title_suffix,
                        "content": chunk
                    })
            else:
                balanced_slides.append(slide)
        
        return balanced_slides
    
    def _create_content_slide(self, prs, slide_content):
        """Create a styled content slide with title and balanced content"""
        bullet_slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Style title
        title_shape = slide.shapes.title
        title_shape.text = slide_content["title"] or "Content"
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(31, 73, 125)  # Dark blue
        
        # Add styled content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.5)
        text_frame.margin_top = Inches(0.3)
        
        for i, content_item in enumerate(slide_content["content"]):
            if i == 0:
                # First content item
                p = text_frame.paragraphs[0]
            else:
                # Additional content items
                p = text_frame.add_paragraph()
            
            # Style content based on type
            if len(content_item) < 100 and not '. ' in content_item:
                # Short bullet point
                p.text = content_item
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(68, 114, 196)  # Medium blue
                p.space_after = Pt(12)
            else:
                # Longer paragraph
                if len(content_item) > 200:
                    # Split long paragraphs
                    sentences = content_item.split('. ')
                    for j, sentence in enumerate(sentences[:3]):
                        if j == 0:
                            p.text = sentence + ('.' if not sentence.endswith('.') else '')
                            p.font.size = Pt(18)
                            p.font.color.rgb = RGBColor(89, 89, 89)  # Dark gray
                            p.space_after = Pt(8)
                        else:
                            p_new = text_frame.add_paragraph()
                            p_new.text = sentence + ('.' if not sentence.endswith('.') else '')
                            p_new.level = 0
                            p_new.font.size = Pt(18)
                            p_new.font.color.rgb = RGBColor(89, 89, 89)
                            p_new.space_after = Pt(8)
                else:
                    p.text = content_item
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(89, 89, 89)  # Dark gray
                    p.space_after = Pt(10)